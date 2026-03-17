"""
Flask web app — Sales Jar Tracker
Upload an Excel file (must contain 'Monthly Target' and 'Current Sales' columns)
and receive a generated jar visualization image.
"""

import base64
import io
import os

import openpyxl
from flask import Flask, render_template, request, jsonify, send_file
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor

from jar_generator import generate_image

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 10 * 1024 * 1024  # 10 MB limit


def _read_excel(file_bytes: bytes) -> tuple[float, float]:
    """
    Parse uploaded Excel bytes.
    Returns (current_sales, monthly_target).
    Raises ValueError with a user-friendly message on bad data.
    """
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    ws = wb.active

    headers = {}
    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, values_only=True):
        for col_idx, cell in enumerate(row):
            if isinstance(cell, str):
                headers[cell.strip()] = col_idx

        if "Monthly Target" in headers and "Current Sales" in headers:
            break

    if "Monthly Target" not in headers or "Current Sales" not in headers:
        raise ValueError(
            "Excel must contain rows with labels 'Monthly Target' and 'Current Sales'."
        )

    target  = None
    current = None

    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, values_only=True):
        if len(row) < 2:
            continue
        label = str(row[0]).strip() if row[0] is not None else ""
        value = row[1]
        if label == "Monthly Target":
            target = value
        elif label == "Current Sales":
            current = value

    if target is None or current is None:
        raise ValueError("Could not find values for 'Monthly Target' and 'Current Sales'.")

    try:
        target  = float(target)
        current = float(current)
    except (TypeError, ValueError):
        raise ValueError("'Monthly Target' and 'Current Sales' values must be numbers.")

    if target <= 0:
        raise ValueError("'Monthly Target' must be greater than 0.")

    return current, target


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/generate", methods=["POST"])
def generate():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded."}), 400

    f = request.files["file"]
    if not f.filename or not f.filename.lower().endswith((".xlsx", ".xlsm")):
        return jsonify({"error": "Please upload an .xlsx or .xlsm Excel file."}), 400

    try:
        file_bytes = f.read()
        current_sales, monthly_target = _read_excel(file_bytes)
        png_bytes = generate_image(current_sales, monthly_target)
    except ValueError as e:
        return jsonify({"error": str(e)}), 422
    except Exception as e:
        return jsonify({"error": f"Unexpected error: {e}"}), 500

    encoded = base64.b64encode(png_bytes).decode("utf-8")
    fill_pct = round(current_sales / monthly_target * 100, 1)

    return jsonify({
        "image":         f"data:image/png;base64,{encoded}",
        "current_sales": current_sales,
        "target":        monthly_target,
        "fill_pct":      fill_pct,
    })


@app.route("/download", methods=["POST"])
def download():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded."}), 400

    f = request.files["file"]
    if not f.filename or not f.filename.lower().endswith((".xlsx", ".xlsm")):
        return jsonify({"error": "Please upload an .xlsx file."}), 400

    try:
        file_bytes = f.read()
        current_sales, monthly_target = _read_excel(file_bytes)
        png_bytes = generate_image(current_sales, monthly_target)
    except ValueError as e:
        return jsonify({"error": str(e)}), 422
    except Exception as e:
        return jsonify({"error": f"Unexpected error: {e}"}), 500

    return send_file(
        io.BytesIO(png_bytes),
        mimetype="image/png",
        as_attachment=True,
        download_name="sales_jar_tracker.png",
    )


def _build_pptx(png_bytes: bytes) -> bytes:
    """
    Wrap a PNG image in a 16:9 PowerPoint slide with a dark-blue background.
    The image is centered and fills the slide with correct aspect ratio.
    """
    prs = Presentation()

    # 16:9 widescreen (13.33 in × 7.5 in)
    slide_w = Inches(13.33)
    slide_h = Inches(7.5)
    prs.slide_width  = slide_w
    prs.slide_height = slide_h

    slide_layout = prs.slide_layouts[6]   # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Dark blue background (matches image gradient mid-point)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(16, 60, 125)

    # Add the PNG image, centred, filling the slide
    img_buf = io.BytesIO(png_bytes)
    pic = slide.shapes.add_picture(img_buf, Emu(0), Emu(0), slide_w, slide_h)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


@app.route("/download_ppt", methods=["POST"])
def download_ppt():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded."}), 400

    f = request.files["file"]
    if not f.filename or not f.filename.lower().endswith((".xlsx", ".xlsm")):
        return jsonify({"error": "Please upload an .xlsx or .xlsm Excel file."}), 400

    try:
        file_bytes = f.read()
        current_sales, monthly_target = _read_excel(file_bytes)
        png_bytes  = generate_image(current_sales, monthly_target)
        pptx_bytes = _build_pptx(png_bytes)
    except ValueError as e:
        return jsonify({"error": str(e)}), 422
    except Exception as e:
        return jsonify({"error": f"Unexpected error: {e}"}), 500

    return send_file(
        io.BytesIO(pptx_bytes),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name="sales_jar_tracker.pptx",
    )


if __name__ == "__main__":
    app.run(debug=True, port=5050)
